#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Inyecta las notas del orador en los 8 pptx, verificando el orden por título."""
import sys
from pptx import Presentation
import notas_b1b2, notas_b3b4b5, notas_b6b7b8

DECKS = {
    "IQS_B1_Introduccion.pptx": notas_b1b2.B1,
    "IQS_B2_Tipologia_Seguridad.pptx": notas_b1b2.B2,
    "IQS_B3_Sensores_Actuadores.pptx": notas_b3b4b5.B3,
    "IQS_B4_Cinematica_Estatica.pptx": notas_b3b4b5.B4,
    "IQS_B5_Modelado_Control.pptx": notas_b3b4b5.B5,
    "IQS_B6_Percepcion_SLAM.pptx": notas_b6b7b8.B6,
    "IQS_B7_ROS2_Planificacion.pptx": notas_b6b7b8.B7,
    "IQS_B8_Robot_Learning.pptx": notas_b6b7b8.B8,
}

fallos = []
for fichero, notas in DECKS.items():
    prs = Presentation(fichero)
    slides = list(prs.slides)
    if len(slides) != len(notas):
        fallos.append(f"{fichero}: {len(slides)} diapositivas pero {len(notas)} notas")
        continue
    for i, (slide, (frag, texto)) in enumerate(zip(slides, notas), 1):
        textos = " | ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        if frag not in textos:
            fallos.append(f"{fichero} diapositiva {i}: no contiene {frag!r}")
            continue
        ns = slide.notes_slide  # crea la parte de notas si no existe
        ns.notes_text_frame.text = texto
    prs.save(fichero)
    print(f"{fichero}: {len(notas)} notas")

if fallos:
    print("FALLOS:")
    for x in fallos:
        print(" -", x)
    sys.exit(1)
print("notas inyectadas")
