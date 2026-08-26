#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Enlaces a Colab en las portadillas de sesión de los decks (post-proceso con python-pptx).

Añade un enlace clicable «Abrir en Colab» en la diapositiva de portadilla de cada
sesión que tiene cuaderno, para poder saltar al cuaderno desde la presentación
proyectada. Es idempotente: si el enlace ya está, no lo duplica.

La correspondencia sesión → cuaderno NO es uno a uno en B6 y B8, porque los
ficheros de los cuadernos de S25-S27 y S39-S40 van desfasados respecto del
calendario del Plan (véase el punto correspondiente del informe de verificación).
Por eso el rótulo lleva siempre el número del cuaderno que abre.

Al reescribir con python-pptx se pierde la entrada Default jpg de [Content_Types].xml
y PowerPoint da el fichero por dañado, así que al terminar se llama a
arreglar_content_types.py, que la repone y es idempotente.

Uso, desde entregables/presentaciones/:
    python3 ../../generadores/enlaces_colab.py [rama]
"""
import glob, os, re, subprocess, sys

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO = "grcarmenaty/MecatronicaRobotica"
RAMA = sys.argv[1] if len(sys.argv) > 1 else "main"
BASE = f"https://colab.research.google.com/github/{REPO}/blob/{RAMA}/cuadernos/"
MARCA = "Abrir en Colab"
VERDE = RGBColor(0x2F, 0xD2, 0x6E)

# deck -> {portadilla de sesión: [(nº de cuaderno, matiz para el rótulo), ...]}
MAPA = {
    "IQS_B3_Sensores_Actuadores.pptx":  {"S8": [("S08", "")], "S9": [("S09", "")], "S10": [("S10", "")],
                                         "S11": [("S11", "")], "S12": [("S12", "")]},
    "IQS_B4_Cinematica_Estatica.pptx":  {"S13": [("S13", "")], "S14": [("S14", "")], "S15": [("S15", "")],
                                         "S16": [("S16", "")], "S17": [("S17", "")], "S18": [("S18", "")],
                                         "S19": [("S19", "")]},
    "IQS_B5_Modelado_Control.pptx":     {"S20": [("S20", "")], "S21": [("S21", "")], "S22": [("S22", "")],
                                         "S23": [("S23", "")], "S24": [("S24", "")]},
    "IQS_B6_Percepcion_SLAM.pptx":      {"S26": [("S25", " imagen y calibración")],
                                         "S27": [("S26", " procesado clásico"),
                                                 ("S27", " percepción aprendida")],
                                         "S28": [("S28", "")], "S29": [("S29", "")], "S30": [("S30", "")],
                                         "S31": [("S31", "")]},
    "IQS_B7_ROS2_Planificacion.pptx":   {"S34": [("S34", "")], "S35": [("S35", "")]},
    "IQS_B8_Robot_Learning.pptx":       {"S38": [("S38", "")],
                                         "S39": [("S39", " RL profundo"),
                                                 ("S40", " imitación y VLA")]},
}

FICHEROS = {os.path.basename(f)[6:9]: os.path.basename(f)
            for f in glob.glob(os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                            "..", "cuadernos", "*.ipynb"))}


def texto(slide):
    return "".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def portadilla(slide):
    """La portadilla empieza por el número de sesión pegado al título («S14Cinemática…»),
    mientras que las de contenido llevan un separador («S14 · DE LAS ARTICULACIONES…»)."""
    t = texto(slide).lstrip()
    m = re.match(r"(S\d{1,2})(?![\d\s·–-])", t)
    return m.group(1) if m else None


def poner_enlace(slide, y, url, rotulo):
    caja = slide.shapes.add_textbox(Inches(2.25), Inches(y), Inches(6.0), Inches(0.38))
    tf = caja.text_frame
    # Anchura fija y sin autoajuste: con wrap="none" + spAutoFit el renderizador
    # encoge la caja al texto manteniendo su centro, y dos rótulos de distinta
    # longitud en la misma portadilla dejan de alinearse por la izquierda.
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r = tf.paragraphs[0].add_run()
    r.text = rotulo
    r.font.name, r.font.size, r.font.bold = "Calibri", Pt(14), True
    r.font.color.rgb = VERDE
    # El enlace se cuelga de la FORMA, no del run. Un hipervínculo de texto hace que
    # LibreOffice (y PowerPoint con el estilo por defecto) repinte el rótulo del azul
    # `hlink`, ilegible sobre el azul marino de la portadilla, e ignore tanto el
    # solidFill del run como el color del tema. Con la acción en la forma, el rótulo
    # conserva el verde de acento y la caja entera queda clicable; LibreOffice la
    # exporta además como anotación /Link, así que el PDF también salta a Colab.
    caja.click_action.hyperlink.address = url


def procesar(ruta):
    prs = Presentation(ruta)
    pendientes = dict(MAPA[os.path.basename(ruta)])
    puestos = 0
    for slide in prs.slides:
        ses = portadilla(slide)
        if ses not in pendientes:
            continue
        if MARCA in texto(slide):          # idempotencia
            pendientes.pop(ses)
            continue
        for i, (num, matiz) in enumerate(pendientes.pop(ses)):
            poner_enlace(slide, 3.62 + 0.40 * i, BASE + FICHEROS[num],
                         f"▸  {MARCA} · cuaderno {num}{matiz}")
            puestos += 1
    if puestos:                 # sin cambios no se reescribe: python-pptx tocaría el zip en vano
        prs.save(ruta)
    estado = "OK" if not pendientes else f"SIN PORTADILLA: {sorted(pendientes)}"
    print(f"{os.path.basename(ruta):34s} {puestos} enlaces  ({estado})")
    return not pendientes, puestos


ok, reescritos = True, []
for nombre in sorted(MAPA):
    if os.path.exists(nombre):
        completo, puestos = procesar(nombre)
        ok &= completo
        if puestos:
            reescritos.append(nombre)
    else:
        print(f"{nombre:34s} NO ENCONTRADO en {os.getcwd()}")
        ok = False
# python-pptx se deja la entrada Default jpg por el camino; el script del repo la repone
ARREGLA = os.path.join(os.path.dirname(os.path.abspath(__file__)), "arreglar_content_types.py")
tocados = reescritos
if tocados:
    print()
    subprocess.run([sys.executable, ARREGLA, *tocados], check=True)

print("\nRama de los enlaces:", RAMA)
sys.exit(0 if ok else 1)
